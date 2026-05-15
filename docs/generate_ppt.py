"""
Generate MarketMind AI Demo PPT
Run: python3 docs/generate_ppt.py
Output: docs/MarketMind_AI_Demo.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colour palette ──────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
CARD_BG      = RGBColor(0x1A, 0x25, 0x3D)   # card navy
ACCENT       = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
GREEN        = RGBColor(0x22, 0xC5, 0x5E)   # green
AMBER        = RGBColor(0xF5, 0x9E, 0x0B)   # amber
RED          = RGBColor(0xEF, 0x44, 0x44)   # red
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY   = RGBColor(0xCB, 0xD5, 0xE1)
DIM          = RGBColor(0x64, 0x74, 0x8B)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── Helpers ─────────────────────────────────────────────────────────────────

def bg(slide, color=DARK_BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=20, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb

def accent_bar(slide, y=Inches(0.08), h=Inches(0.06)):
    """Thin accent line at top of slide."""
    box(slide, 0, y, W, h, ACCENT)

def slide_number(slide, n):
    txt(slide, str(n), Inches(12.8), Inches(7.1), Inches(0.4), Inches(0.3),
        size=10, color=DIM, align=PP_ALIGN.RIGHT)

def tag_pill(slide, label, x, y, color=ACCENT):
    """Small coloured pill label."""
    box(slide, x, y, Inches(1.6), Inches(0.32), color)
    txt(slide, label, x, y, Inches(1.6), Inches(0.32),
        size=10, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

def kpi_card(slide, title, value, sub, x, y, color=ACCENT):
    box(slide, x, y, Inches(2.6), Inches(1.5), CARD_BG)
    # left colour stripe
    box(slide, x, y, Inches(0.06), Inches(1.5), color)
    txt(slide, title, x+Inches(0.12), y+Inches(0.1), Inches(2.4), Inches(0.35),
        size=11, color=LIGHT_GREY)
    txt(slide, value, x+Inches(0.12), y+Inches(0.42), Inches(2.4), Inches(0.55),
        size=26, bold=True, color=color)
    txt(slide, sub, x+Inches(0.12), y+Inches(0.98), Inches(2.4), Inches(0.35),
        size=10, color=DIM)

def bullet_card(slide, title, items, x, y, w=Inches(5.8), color=ACCENT):
    bh = Inches(0.42 + 0.38 * len(items) + 0.15)
    box(slide, x, y, w, bh, CARD_BG)
    box(slide, x, y, Inches(0.06), bh, color)
    txt(slide, title, x+Inches(0.16), y+Inches(0.1), w-Inches(0.3), Inches(0.32),
        size=13, bold=True, color=color)
    for i, item in enumerate(items):
        txt(slide, f"• {item}", x+Inches(0.18), y+Inches(0.42)+Inches(0.38*i),
            w-Inches(0.35), Inches(0.35), size=12, color=LIGHT_GREY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

# Gradient-like diagonal shape
shape = s.shapes.add_shape(1, Inches(7.5), Inches(0), Inches(6), Inches(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = CARD_BG
shape.line.fill.background()

# Brand name
txt(s, "MarketMind AI", Inches(0.6), Inches(1.6), Inches(7), Inches(1.2),
    size=52, bold=True, color=WHITE)

# Accent underline
box(s, Inches(0.6), Inches(2.85), Inches(3.2), Inches(0.07), ACCENT)

txt(s, "Autonomous Multi-Agent Financial Intelligence", Inches(0.6), Inches(3.0),
    Inches(7), Inches(0.6), size=20, color=LIGHT_GREY, italic=True)

txt(s, "MVP Demo  •  May 2026", Inches(0.6), Inches(3.72), Inches(5), Inches(0.4),
    size=14, color=DIM)

# Right panel text
txt(s, "🤖  6 AI Agents", Inches(8.1), Inches(1.8), Inches(4.5), Inches(0.5),
    size=18, color=ACCENT)
txt(s, "📈  Real-time NSE Data", Inches(8.1), Inches(2.4), Inches(4.5), Inches(0.5),
    size=18, color=WHITE)
txt(s, "🧠  Gemini 2.5-Flash LLM", Inches(8.1), Inches(3.0), Inches(4.5), Inches(0.5),
    size=18, color=WHITE)
txt(s, "📧  Personalised Email Digests", Inches(8.1), Inches(3.6), Inches(4.5), Inches(0.5),
    size=18, color=WHITE)
txt(s, "🌐  India + Global Market News", Inches(8.1), Inches(4.2), Inches(4.5), Inches(0.5),
    size=18, color=WHITE)
txt(s, "⏰  Fully Automated  •  Zero Manual Work", Inches(8.1), Inches(4.8),
    Inches(4.5), Inches(0.5), size=18, color=WHITE)

slide_number(s, 1)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem We're Solving
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
txt(s, "The Problem", Inches(0.6), Inches(0.3), Inches(8), Inches(0.6),
    size=32, bold=True, color=WHITE)
box(s, Inches(0.6), Inches(0.95), Inches(1.8), Inches(0.05), ACCENT)

problems = [
    ("🔴  Information Overload",   "100s of news articles daily — impossible to track every stock manually"),
    ("🔴  Missed Opportunities",   "Market moves happen fast; retail investors react too late"),
    ("🔴  No Personalisation",     "Generic newsletters don't match your watchlist or risk appetite"),
    ("🔴  Manual Analysis",        "Chart reading, sentiment tracking — time-consuming and error-prone"),
    ("🔴  Fragmented Data",        "Price data, news, and AI analysis live in different places"),
]

for i, (title, desc) in enumerate(problems):
    y = Inches(1.2) + Inches(1.1 * i)
    box(s, Inches(0.55), y, Inches(12.2), Inches(0.95), CARD_BG)
    box(s, Inches(0.55), y, Inches(0.07), Inches(0.95), RED)
    txt(s, title, Inches(0.78), y+Inches(0.08), Inches(4), Inches(0.4),
        size=14, bold=True, color=WHITE)
    txt(s, desc, Inches(4.9), y+Inches(0.08), Inches(7.7), Inches(0.75),
        size=13, color=LIGHT_GREY)

slide_number(s, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Our Solution
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
txt(s, "Our Solution", Inches(0.6), Inches(0.3), Inches(8), Inches(0.6),
    size=32, bold=True, color=WHITE)
box(s, Inches(0.6), Inches(0.95), Inches(1.8), Inches(0.05), ACCENT)

txt(s, "MarketMind AI is a fully autonomous platform that monitors NSE markets,\nruns AI analysis, and delivers personalised signals to your inbox — 24/7.",
    Inches(0.6), Inches(1.1), Inches(12), Inches(0.8), size=15, color=LIGHT_GREY)

solutions = [
    (GREEN,  "✅  Zero Manual Work",          "Pipeline runs on schedule — pre-market 08:45 & post-market 16:15 IST"),
    (GREEN,  "✅  AI-Powered Signals",         "Gemini 2.5-Flash analyses each stock and generates BUY / SELL / HOLD"),
    (GREEN,  "✅  Personalised Digests",       "Each subscriber gets their own watchlist digest, not a generic blast"),
    (GREEN,  "✅  India + Global News",        "NSE stocks news + Reuters, Yahoo Finance, Google News"),
    (GREEN,  "✅  One Dashboard",              "Signals, news, analysis, alerts — all in one real-time web UI"),
]

for i, (color, title, desc) in enumerate(solutions):
    y = Inches(2.05) + Inches(1.0 * i)
    box(s, Inches(0.55), y, Inches(12.2), Inches(0.88), CARD_BG)
    box(s, Inches(0.55), y, Inches(0.07), Inches(0.88), color)
    txt(s, title, Inches(0.78), y+Inches(0.08), Inches(4), Inches(0.35),
        size=14, bold=True, color=WHITE)
    txt(s, desc, Inches(4.9), y+Inches(0.08), Inches(7.7), Inches(0.68),
        size=13, color=LIGHT_GREY)

slide_number(s, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
txt(s, "System Architecture", Inches(0.6), Inches(0.3), Inches(10), Inches(0.6),
    size=32, bold=True, color=WHITE)
box(s, Inches(0.6), Inches(0.95), Inches(2.4), Inches(0.05), ACCENT)

agents = [
    (ACCENT,  "1  MarketDataAgent",     "NSE live price, volume,\nP/E, 52w high/low"),
    (ACCENT,  "2  NewsAgent",           "Indian + Global RSS feeds,\nsentiment classification"),
    (AMBER,   "3  AIAnalysisAgent",     "Gemini 2.5-Flash\nBullish / Bearish / Neutral"),
    (GREEN,   "4  SignalGenerator",     "BUY / SELL / HOLD\nconfidence + risk level"),
    (ACCENT,  "5  ReportGenerator",     "Charts + HTML\nsummary report"),
    (ACCENT,  "6  EmailAlertAgent",     "Personalised digest\n+ unsubscribe footer"),
]

ax = Inches(0.35)
for i, (color, name, desc) in enumerate(agents):
    x = ax + Inches(2.15 * i)
    box(s, x, Inches(1.25), Inches(1.95), Inches(2.2), CARD_BG)
    box(s, x, Inches(1.25), Inches(1.95), Inches(0.08), color)
    txt(s, name, x+Inches(0.1), Inches(1.38), Inches(1.8), Inches(0.45),
        size=11, bold=True, color=color)
    txt(s, desc, x+Inches(0.1), Inches(1.85), Inches(1.78), Inches(0.7),
        size=10, color=LIGHT_GREY)
    # Arrow (except last)
    if i < 5:
        txt(s, "→", x+Inches(1.95), Inches(1.9), Inches(0.25), Inches(0.35),
            size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Orchestrator bar
box(s, Inches(0.35), Inches(3.65), Inches(12.6), Inches(0.45), RGBColor(0x1E, 0x3A, 0x5F))
txt(s, "Orchestrator  (src/orchestrator.py)  —  coordinates full pipeline, applies per-subscriber overrides",
    Inches(0.5), Inches(3.7), Inches(12.2), Inches(0.38),
    size=12, color=ACCENT, align=PP_ALIGN.CENTER)

# Stack row
stack = [
    ("Python 3.10", ACCENT), ("Flask 3.x", WHITE), ("SQLite / SQLAlchemy", WHITE),
    ("APScheduler", AMBER), ("Chart.js", WHITE), ("Bootstrap 5", WHITE), ("Gmail SMTP", GREEN),
]
sx = Inches(0.5)
for label, color in stack:
    w = Inches(1.55)
    box(s, sx, Inches(4.3), w, Inches(0.38), CARD_BG)
    txt(s, label, sx, Inches(4.3), w, Inches(0.38),
        size=10, bold=True, color=color, align=PP_ALIGN.CENTER)
    sx += Inches(1.7)

txt(s, "Tech Stack", Inches(0.5), Inches(4.1), Inches(4), Inches(0.28),
    size=11, color=DIM)

# DB
box(s, Inches(0.35), Inches(4.9), Inches(12.6), Inches(1.8), CARD_BG)
box(s, Inches(0.35), Inches(4.9), Inches(12.6), Inches(0.07), DIM)
txt(s, "SQLite Database  —  6 tables: StockData · NewsRecord · AnalysisReport · TradingSignal · AlertRecord · SubscriberRecord",
    Inches(0.6), Inches(5.0), Inches(12), Inches(0.38), size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
txt(s, "REST API  —  19 endpoints  ·  Flask-served Jinja2 HTML  ·  JSON APIs consumed by Chart.js + vanilla JS",
    Inches(0.6), Inches(5.45), Inches(12), Inches(0.38), size=11, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
txt(s, "Scheduled: APScheduler CronTrigger  08:45 IST (pre-market)  +  16:15 IST (post-market)",
    Inches(0.6), Inches(5.9), Inches(12), Inches(0.38), size=11, color=AMBER, align=PP_ALIGN.CENTER)

slide_number(s, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Key Features Achieved (MVP)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
txt(s, "What We've Built  —  MVP Features", Inches(0.6), Inches(0.3),
    Inches(10), Inches(0.6), size=32, bold=True, color=WHITE)
box(s, Inches(0.6), Inches(0.95), Inches(3.5), Inches(0.05), ACCENT)

col1 = [
    ("Real-time Dashboard",       "Signals, news, charts, alerts in one UI"),
    ("6-Agent AI Pipeline",       "End-to-end from raw data → email digest"),
    ("Gemini 2.5-Flash Analysis", "Per-stock sentiment + key factors from LLM"),
    ("BUY / SELL / HOLD Signals", "Confidence score + risk level per signal"),
    ("1700+ NSE Stocks",          "Live NSE CSV catalog, 6-hour cache"),
    ("India + Global News",       "6 international RSS feeds + per-stock feeds"),
]
col2 = [
    ("Email Subscription System", "Self-service subscribe / update / unsubscribe"),
    ("Personalised Digests",      "Each subscriber gets their own watchlist email"),
    ("Welcome & Update Emails",   "Sent immediately on subscribe / watchlist change"),
    ("One-Click Unsubscribe",     "UUID token in every email footer"),
    ("Automated Scheduling",      "08:45 pre-market + 16:15 post-market, daily"),
    ("Responsive Mobile UI",      "Dark/light theme, scrollable panels, bottom nav"),
]

for i, (feat, desc) in enumerate(col1):
    y = Inches(1.2) + Inches(0.95 * i)
    box(s, Inches(0.4), y, Inches(6.1), Inches(0.82), CARD_BG)
    box(s, Inches(0.4), y, Inches(0.07), Inches(0.82), GREEN)
    txt(s, feat, Inches(0.6), y+Inches(0.06), Inches(3), Inches(0.32),
        size=12, bold=True, color=GREEN)
    txt(s, desc, Inches(0.6), y+Inches(0.4), Inches(5.7), Inches(0.32),
        size=11, color=LIGHT_GREY)

for i, (feat, desc) in enumerate(col2):
    y = Inches(1.2) + Inches(0.95 * i)
    box(s, Inches(6.85), y, Inches(6.1), Inches(0.82), CARD_BG)
    box(s, Inches(6.85), y, Inches(0.07), Inches(0.82), ACCENT)
    txt(s, feat, Inches(7.05), y+Inches(0.06), Inches(3.2), Inches(0.32),
        size=12, bold=True, color=ACCENT)
    txt(s, desc, Inches(7.05), y+Inches(0.4), Inches(5.7), Inches(0.32),
        size=11, color=LIGHT_GREY)

slide_number(s, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Live Demo Walkthrough
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
txt(s, "Live Demo  —  Walkthrough", Inches(0.6), Inches(0.3),
    Inches(10), Inches(0.6), size=32, bold=True, color=WHITE)
box(s, Inches(0.6), Inches(0.95), Inches(2.6), Inches(0.05), ACCENT)

steps = [
    (ACCENT,  "Step 1  —  Dashboard",
              "KPIs · scrollable signals · news tabs (India / Global) · recent alerts"),
    (ACCENT,  "Step 2  —  Run Analysis",
              "Stock picker modal → select stocks → trigger 6-agent pipeline → live status polling"),
    (GREEN,   "Step 3  —  Stock Detail",
              "Price chart · signal history · Gemini analysis · news feed per symbol"),
    (GREEN,   "Step 4  —  Subscribe",
              "Enter email · pick stocks by sector · subscribe → welcome email in inbox"),
    (AMBER,   "Step 5  —  Update Watchlist",
              "Re-submit same email with different stocks → update email confirming change"),
    (AMBER,   "Step 6  —  Unsubscribe",
              "Manage subscription card → lookup email → click link in email → confirmation"),
    (RED,     "Step 7  —  Alert History",
              "Full digest history · sent/failed status · recipient list"),
]

for i, (color, title, desc) in enumerate(steps):
    y = Inches(1.2) + Inches(0.87 * i)
    box(s, Inches(0.4), y, Inches(12.5), Inches(0.75), CARD_BG)
    box(s, Inches(0.4), y, Inches(0.07), Inches(0.75), color)
    # Step number bubble
    nb = s.shapes.add_shape(9, Inches(0.55), y+Inches(0.12),
                             Inches(0.5), Inches(0.5))
    nb.fill.solid(); nb.fill.fore_color.rgb = color
    nb.line.fill.background()
    txt(s, str(i+1), Inches(0.55), y+Inches(0.12), Inches(0.5), Inches(0.5),
        size=12, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    txt(s, title, Inches(1.2), y+Inches(0.06), Inches(3.8), Inches(0.3),
        size=12, bold=True, color=color)
    txt(s, desc, Inches(5.2), y+Inches(0.06), Inches(7.6), Inches(0.6),
        size=11, color=LIGHT_GREY)

slide_number(s, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Metrics / KPIs
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
txt(s, "Platform at a Glance", Inches(0.6), Inches(0.3),
    Inches(10), Inches(0.6), size=32, bold=True, color=WHITE)
box(s, Inches(0.6), Inches(0.95), Inches(2.5), Inches(0.05), ACCENT)

kpis = [
    ("Stocks Supported",  "1,700+",  "Live NSE catalog",          ACCENT),
    ("AI Agents",         "6",       "End-to-end pipeline",        GREEN),
    ("REST Endpoints",    "19",      "Full API coverage",          AMBER),
    ("News Sources",      "10+",     "India + 6 global RSS",       ACCENT),
    ("Email Events",      "3",       "Welcome · Update · Digest",  GREEN),
    ("DB Tables",         "6",       "SQLite / SQLAlchemy ORM",    AMBER),
    ("Digests Per Day",   "2",       "08:45 & 16:15 IST",          RED),
    ("Setup Time",        "< 5 min", "Install → running",          ACCENT),
]

cols = 4
for i, (title, value, sub, color) in enumerate(kpis):
    row = i // cols
    col = i % cols
    x = Inches(0.5) + Inches(3.2 * col)
    y = Inches(1.3) + Inches(2.0 * row)
    kpi_card(s, title, value, sub, x, y, color)

slide_number(s, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Roadmap
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
txt(s, "Roadmap  —  What's Next", Inches(0.6), Inches(0.3),
    Inches(10), Inches(0.6), size=32, bold=True, color=WHITE)
box(s, Inches(0.6), Inches(0.95), Inches(2.4), Inches(0.05), ACCENT)

phases = [
    ("Phase 2  —  Near Term\n(1–2 months)", ACCENT, [
        "Technical indicators: RSI, MACD, Bollinger Bands",
        "Signal backtesting & P&L performance tracking",
        "Telegram / WhatsApp alert notifications",
        "Stock screener with custom filter criteria",
        "User login & per-user portfolio management",
    ]),
    ("Phase 3  —  Growth\n(3–6 months)", AMBER, [
        "PostgreSQL migration for production scale",
        "Paid data API (Zerodha Kite / Upstox) for reliability",
        "SendGrid / AWS SES for email scale",
        "Mobile app (React Native) — push notifications",
        "Multi-market support: BSE, global indices",
    ]),
    ("Phase 4  —  Scale\n(6–12 months)", GREEN, [
        "Options & derivatives signal support",
        "Social sentiment: Twitter/X + Reddit scraping",
        "Tiered SaaS subscription model",
        "Broker API integration — one-click trade execution",
        "Regulatory compliance & SEBI advisory framework",
    ]),
]

for i, (phase, color, items) in enumerate(phases):
    x = Inches(0.4) + Inches(4.3 * i)
    box(s, x, Inches(1.2), Inches(4.05), Inches(5.8), CARD_BG)
    box(s, x, Inches(1.2), Inches(4.05), Inches(0.08), color)
    txt(s, phase, x+Inches(0.18), Inches(1.28), Inches(3.7), Inches(0.7),
        size=13, bold=True, color=color)
    for j, item in enumerate(items):
        txt(s, f"→  {item}",
            x+Inches(0.18), Inches(2.05)+Inches(0.88*j),
            Inches(3.7), Inches(0.78), size=11, color=LIGHT_GREY)

slide_number(s, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Why MarketMind AI Wins
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)
txt(s, "Why MarketMind AI?", Inches(0.6), Inches(0.3),
    Inches(10), Inches(0.6), size=32, bold=True, color=WHITE)
box(s, Inches(0.6), Inches(0.95), Inches(2.0), Inches(0.05), ACCENT)

comparisons = [
    ("",                    "MarketMind AI",  "Generic Newsletter",  "Manual Research"),
    ("AI Analysis",         "✅ Per stock",    "❌ Generic",           "❌ Manual effort"),
    ("Personalisation",     "✅ Per watchlist","❌ Same for all",      "✅ Manual only"),
    ("Automation",          "✅ Fully auto",   "⚠️ Semi-auto",         "❌ None"),
    ("Frequency",           "✅ 2× per day",   "⚠️ Daily/weekly",      "❌ Ad-hoc"),
    ("Real-time Data",      "✅ NSE live",     "⚠️ Delayed",           "⚠️ Depends"),
    ("Global News",         "✅ 6 feeds",      "⚠️ India only",        "❌ Manual"),
    ("Unsubscribe Control", "✅ One-click",    "⚠️ Email footer",      "N/A"),
    ("Cost",                "✅ Open source",  "💰 Paid",              "⏰ Your time"),
]

col_widths = [Inches(2.8), Inches(2.8), Inches(3.2), Inches(3.2)]
col_xs     = [Inches(0.4), Inches(3.3), Inches(6.2), Inches(9.5)]
row_h = Inches(0.6)

for row_i, row in enumerate(comparisons):
    y = Inches(1.1) + Inches(row_h * row_i)
    for col_i, cell in enumerate(row):
        bg_c = CARD_BG if row_i > 0 else RGBColor(0x0A, 0x12, 0x22)
        if col_i == 1 and row_i > 0:
            bg_c = RGBColor(0x0C, 0x2A, 0x1A)  # green tint for our column
        box(s, col_xs[col_i], y, col_widths[col_i]-Inches(0.05), row_h-Inches(0.04), bg_c)
        fc = ACCENT if row_i == 0 else (GREEN if col_i == 1 else LIGHT_GREY)
        if col_i == 1 and row_i == 0:
            fc = GREEN
        txt(s, cell, col_xs[col_i]+Inches(0.1), y+Inches(0.1),
            col_widths[col_i]-Inches(0.2), row_h-Inches(0.15),
            size=12 if row_i > 0 else 13,
            bold=(row_i == 0 or col_i == 0),
            color=fc, align=PP_ALIGN.CENTER if col_i > 0 else PP_ALIGN.LEFT)

slide_number(s, 9)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Thank You / Q&A
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
accent_bar(s)

# Big centre block
box(s, Inches(1.5), Inches(1.2), Inches(10.3), Inches(4.5), CARD_BG)
box(s, Inches(1.5), Inches(1.2), Inches(10.3), Inches(0.1), ACCENT)

txt(s, "Thank You", Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.2),
    size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, Inches(4.5), Inches(2.82), Inches(4.3), Inches(0.07), ACCENT)

txt(s, "Questions & Demo", Inches(1.5), Inches(2.95), Inches(10.3), Inches(0.6),
    size=22, color=LIGHT_GREY, align=PP_ALIGN.CENTER, italic=True)

txt(s, "MarketMind AI  •  MVP v1.0  •  May 2026",
    Inches(1.5), Inches(3.7), Inches(10.3), Inches(0.5),
    size=14, color=DIM, align=PP_ALIGN.CENTER)

txt(s, "🌐  http://localhost:5050",
    Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.5),
    size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

txt(s, "Stack:  Python 3.10  ·  Flask  ·  SQLite  ·  Gemini 2.5-Flash  ·  APScheduler  ·  Gmail SMTP",
    Inches(1.5), Inches(5.05), Inches(10.3), Inches(0.4),
    size=12, color=DIM, align=PP_ALIGN.CENTER)

slide_number(s, 10)


# ── Save ────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "MarketMind_AI_Demo.pptx")
prs.save(out)
print(f"✅  Saved → {out}")
